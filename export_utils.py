"""
export_utils.py — Geração de relatórios exportáveis em Excel, PDF e
PowerPoint a partir de DataFrames (e, opcionalmente, imagens de gráficos
Plotly já renderizados em PNG).

Dependências: openpyxl (Excel), reportlab (PDF), python-pptx (PowerPoint).
Todas listadas em requirements.txt. Se alguma não estiver disponível no
ambiente, a função correspondente levanta um erro amigável capturado pela
tela do Streamlit (ver app.py) em vez de derrubar o app inteiro.
"""

import io
import pandas as pd


# ----------------------------------------------------------------------
# Excel
# ----------------------------------------------------------------------
def to_excel_bytes(planilhas: dict) -> bytes:
    """planilhas: {nome_aba: DataFrame}. Retorna bytes de um .xlsx com 1 aba por item."""
    buf = io.BytesIO()
    with pd.ExcelWriter(buf, engine="openpyxl") as writer:
        for nome, df in planilhas.items():
            nome_aba = str(nome)[:31] or "Sheet1"
            (df if isinstance(df, pd.DataFrame) else pd.DataFrame(df)).to_excel(writer, sheet_name=nome_aba, index=False)
    return buf.getvalue()


# ----------------------------------------------------------------------
# Gráfico Plotly -> PNG (opcional; requer kaleido)
# ----------------------------------------------------------------------
def fig_to_png_bytes(fig):
    try:
        return fig.to_image(format="png", width=1000, height=550, scale=2)
    except Exception:
        return None


# ----------------------------------------------------------------------
# PDF (reportlab)
# ----------------------------------------------------------------------
def to_pdf_bytes(titulo: str, secoes: list) -> bytes:
    """
    secoes: lista de tuplas (subtitulo:str, df:DataFrame|None, imagem_png_bytes:bytes|None, texto:str|None)
    Gera um PDF simples: título, e para cada seção um subtítulo + tabela (se houver)
    + imagem (se houver) + texto livre (se houver).
    """
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.lib import colors
    from reportlab.lib.units import cm
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, Image
    from reportlab.lib.styles import getSampleStyleSheet

    buf = io.BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=landscape(A4), topMargin=1.5 * cm, bottomMargin=1.5 * cm)
    styles = getSampleStyleSheet()
    story = [Paragraph(titulo, styles["Title"]), Spacer(1, 0.5 * cm)]

    for subtitulo, df, imagem_bytes, texto in secoes:
        if subtitulo:
            story.append(Paragraph(subtitulo, styles["Heading2"]))
        if texto:
            story.append(Paragraph(texto, styles["Normal"]))
            story.append(Spacer(1, 0.2 * cm))
        if imagem_bytes:
            story.append(Image(io.BytesIO(imagem_bytes), width=22 * cm, height=11 * cm))
            story.append(Spacer(1, 0.3 * cm))
        if df is not None and not df.empty:
            df_show = df.copy()
            for c in df_show.columns:
                df_show[c] = df_show[c].astype(str)
            data = [list(df_show.columns)] + df_show.values.tolist()
            # limita a 40 linhas por tabela para não estourar o PDF
            if len(data) > 41:
                data = data[:41]
            t = Table(data, repeatRows=1)
            t.setStyle(TableStyle([
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#6a1b9a")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTSIZE", (0, 0), (-1, -1), 7),
                ("GRID", (0, 0), (-1, -1), 0.25, colors.grey),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f2f2f2")]),
            ]))
            story.append(t)
        story.append(Spacer(1, 0.6 * cm))

    doc.build(story)
    return buf.getvalue()


# ----------------------------------------------------------------------
# PowerPoint (python-pptx)
# ----------------------------------------------------------------------
def to_pptx_bytes(titulo: str, secoes: list) -> bytes:
    """
    secoes: lista de tuplas (subtitulo:str, df:DataFrame|None, imagem_png_bytes:bytes|None, texto:str|None)
    Gera 1 slide de capa + 1 slide por seção (imagem OU tabela resumida + texto).
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # slide de capa
    capa = prs.slides.add_slide(prs.slide_layouts[0])
    capa.shapes.title.text = titulo
    if len(capa.placeholders) > 1:
        capa.placeholders[1].text = "Sistema de Inteligência de Ativos — Metodologia DRC"

    for subtitulo, df, imagem_bytes, texto in secoes:
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # título em branco
        if slide.shapes.title:
            slide.shapes.title.text = subtitulo or ""

        top = Inches(1.3)
        if texto:
            txbox = slide.shapes.add_textbox(Inches(0.5), top, Inches(12.3), Inches(0.8))
            tf = txbox.text_frame
            tf.text = texto
            for p in tf.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(14)
            top = Inches(2.1)

        if imagem_bytes:
            slide.shapes.add_picture(io.BytesIO(imagem_bytes), Inches(0.7), top, width=Inches(11.8))
        elif df is not None and not df.empty:
            df_show = df.head(15).copy()
            for c in df_show.columns:
                df_show[c] = df_show[c].astype(str)
            rows, cols = df_show.shape[0] + 1, df_show.shape[1]
            table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), top, Inches(12.3), Inches(4.8))
            table = table_shape.table
            for j, col in enumerate(df_show.columns):
                table.cell(0, j).text = str(col)
            for i in range(df_show.shape[0]):
                for j in range(cols):
                    table.cell(i + 1, j).text = str(df_show.iat[i, j])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
